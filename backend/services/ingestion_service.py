"""File ingestion background pipeline (Spec §4.2, §3.1).

run(document_id) is invoked by the worker process (services.worker) after it
claims a pending document from the queue. Opens its own SessionLocal because
no request-scoped DB session exists in the worker.

Pipeline:
  1. Load Document; load the blob via services.object_store (R2 in prod,
     local disk in dev), not settings.uploads_path directly.
  2. _extract(blob, filename) -> [(page_num | None, text), ...] by extension.
     - .pdf  : pypdf -> [(page_num, text), ...]
     - .pptx : python-pptx -> [(slide_num, text), ...]
     - .txt / .md / .markdown : [(None, full_text)]
  3. lib.chunking.chunk_text (500 / 50 overlap).
  4. _embed_and_store: per EMBED_BATCH slice (100 chunks) -- cap-check,
     embed (litellm.embedding), insert (pgvector_store.insert_chunks),
     meter (cost_meter.meter_embedding_response), commit. No DB transaction
     is held open during the litellm.embedding HTTP call (F-02/B-02), and
     nothing accumulates for the whole document in memory at once (F-03) --
     a batch's vectors and raw response are discarded as soon as that
     batch's commit lands.
  5. lib.keyword_index.merge_into_session(stems).
  6. Document.status = ready, page_count populated (None for plaintext).

Because embedding/insert/meter now commit per batch, a crash mid-document
leaves a prefix of chunks durably persisted with their spend already on the
ledger -- nothing to re-record on failure. Resume contract: _embed_and_store
looks up chunk_index values already persisted for the document
(_existing_chunk_indexes) and skips them before making any embedding call,
so re-running ingestion for a partially-ingested document never re-pays for
chunks it already bought and stored.

Steps 5-6 still share one transaction with the final batch's own commit
boundary: merge_into_session only flushes, and the closing db.commit() after
step 6 covers it (F-27 applies only to this tail, not to the already-durable
embedding batches). On exception at any step: db.rollback() first (discards
any unflushed/uncommitted work from this run -- committed batches are
unaffected), then status=failed, error=str(exc)[:1000], committed alone.
"""

import io
import logging
import os

import litellm
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import select

from config import settings
from db.database import SessionLocal
from db.models import ChunkEmbedding, Document
from db.models import Session as SessionModel
from lib import chunking, keyword_index, llm_retry
from services import cost_meter, object_store, pgvector_store


log = logging.getLogger(__name__)

EMBED_BATCH = 100


def _load_blob(store: "object_store.ObjectStore", doc: Document) -> bytes:
    try:
        return store.get(object_store.key_for(doc.id, doc.filename))
    except object_store.ObjectNotFound:
        # Not under the canonical key - try the legacy layout below.
        log.debug("blob miss on canonical key for doc %s; trying legacy key", doc.id)
    # Legacy fallback: pre-F-15 uploads were stored under the bare filename.
    # LocalDiskStore path-containment rejects traversal in doc.filename.
    try:
        return store.get(doc.filename)
    except object_store.ObjectNotFound:
        raise RuntimeError(f"uploaded file not found in object store: {doc.filename}") from None


def _extract_pages(blob: bytes) -> list[tuple[int, str]]:
    reader = PdfReader(io.BytesIO(blob))
    return [(i + 1, (page.extract_text() or "")) for i, page in enumerate(reader.pages)]


def _extract_slides(blob: bytes) -> list[tuple[int, str]]:
    prs = Presentation(io.BytesIO(blob))
    out: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        out.append((i, "\n".join(p for p in parts if p)))
    return out


def _extract_plaintext(blob: bytes) -> list[tuple[None, str]]:
    return [(None, blob.decode("utf-8", errors="replace"))]


def _extract(blob: bytes, filename: str) -> list[tuple[int | None, str]]:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return _extract_pages(blob)
    if ext == ".pptx":
        return _extract_slides(blob)
    if ext in (".txt", ".md", ".markdown"):
        return _extract_plaintext(blob)
    raise ValueError(f"unsupported file type: {ext!r}")


def _existing_chunk_indexes(db, document_id: int) -> set[int]:
    """Chunk indexes already persisted for this document (resume skip-set)."""
    return {
        idx
        for (idx,) in db.execute(
            select(ChunkEmbedding.chunk_index).where(
                ChunkEmbedding.document_id == document_id
            )
        )
    }


def _embed_and_store(db, doc, chunks, *, user_id: str | None) -> int:
    """Embed and persist chunks in EMBED_BATCH slices.

    Per slice: cap-check -> embed (no open transaction during the HTTP
    call) -> insert -> meter -> commit. Skips chunk indexes already
    persisted by a previous attempt, so re-runs never re-pay.
    """
    existing = _existing_chunk_indexes(db, doc.id)
    todo = [c for c in chunks if c.chunk_idx not in existing]
    stored = 0
    for i in range(0, len(todo), EMBED_BATCH):
        batch = todo[i : i + EMBED_BATCH]
        if user_id is not None:
            cost_meter.assert_within_caps(db, user_id)
        db.commit()  # release the connection before the network call
        try:
            resp = llm_retry.retry_sync(
                lambda: litellm.embedding(
                    model=settings.embedding_model,
                    input=[c.text for c in batch],
                    dimensions=settings.embedding_dim,
                    timeout=settings.embedding_timeout_s,
                )
            )
        except Exception as e:
            raise RuntimeError(f"embedding api failed: {e}") from e
        rows = [
            (
                c.chunk_idx,
                c.page,
                c.text,
                item["embedding"] if isinstance(item, dict) else item.embedding,
            )
            for c, item in zip(batch, resp.data)
        ]
        stored += pgvector_store.insert_chunks(
            db, session_id=doc.session_id, document_id=doc.id, rows=rows
        )
        # F-19: ingestion is the largest embedding spender; meter it. The cap
        # gate itself runs at the top of the loop above, before each batch's
        # litellm.embedding call (audit B-01).
        cost_meter.meter_embedding_response(
            db, resp, user_id=user_id, session_id=doc.session_id,
            texts=[c.text for c in batch],
        )
        db.commit()
    return stored


def run(document_id: int) -> None:
    db = SessionLocal()
    try:
        doc = db.get(Document, document_id)
        if doc is None:
            log.warning("ingestion run: document %s not found", document_id)
            return

        owner_id: str | None = None
        # Hoisted local: except arms below read this after db.rollback(),
        # which expires ORM instances -- reading doc.session_id there would
        # trigger an implicit refresh SELECT that can itself raise if the
        # original failure was DB/connection-related, skipping the
        # mark-failed commit and stranding the doc in "pending".
        session_id = doc.session_id
        log.info(
            "ingestion start document_id=%s session_id=%s filename=%s",
            doc.id, session_id, doc.filename,
        )
        try:
            blob = _load_blob(object_store.get_store(), doc)
            pages = _extract(blob, doc.filename)
            # Count only pages/slides that carry a page number; plaintext yields
            # (None, text) so its sum is 0, which we collapse to None (no page
            # concept). Degenerate inputs (0-slide pptx) likewise store None.
            doc.page_count = sum(1 for p, _ in pages if p is not None) or None

            chunks = chunking.chunk_text(pages)
            if not chunks:
                doc.status = "ready"
                db.commit()
                return

            if len(chunks) > settings.max_chunks:
                doc.status = "failed"
                doc.error = "document too large to ingest (chunk limit)"
                db.commit()
                return

            owner_id = db.execute(
                select(SessionModel.user_id).where(SessionModel.id == doc.session_id)
            ).scalar_one_or_none()
            _embed_and_store(db, doc, chunks, user_id=owner_id)

            stems: set[str] = set()
            for c in chunks:
                stems |= keyword_index.build_from_text(c.text)
            if stems:
                keyword_index.merge_into_session(db, doc.session_id, stems)

            log.info(
                "ingestion done document_id=%s session_id=%s chunks=%s",
                doc.id, doc.session_id, len(chunks),
            )
            doc.status = "ready"
            doc.error = None
            db.commit()
        except cost_meter.CostCapExceeded:
            db.rollback()
            log.warning(
                "ingestion stopped: cost cap reached "
                "document_id=%s session_id=%s error=%s",
                document_id, session_id, "cost_cap_exceeded",
                extra={"doc_id": document_id},
            )
            # Per-batch commits mean any batches embedded before the breach
            # are already durable, with their spend already on the ledger --
            # nothing to re-record here.
            doc = db.get(Document, document_id)
            if doc is not None:
                doc.status = "failed"
                doc.error = "daily cost cap reached; ingestion stopped"
                db.commit()
            return
        except Exception as e:
            db.rollback()
            log.error(
                "ingestion failed document_id=%s session_id=%s error=%s",
                document_id, session_id, e,
                extra={"err_type": type(e).__name__, "doc_id": document_id},
                exc_info=settings.env != "prod",
            )
            # Per-batch commits (F-02/F-03/B-02) mean any chunks embedded,
            # inserted, and metered before this failure are already durable
            # on a prior commit -- db.rollback() above only discards this
            # attempt's uncommitted tail (e.g. page_count, keyword merge),
            # never a paid-for batch. No re-record needed.
            doc = db.get(Document, document_id)
            if doc is not None:
                doc.status = "failed"
                doc.error = str(e)[:1000]
                db.commit()
    finally:
        db.close()
